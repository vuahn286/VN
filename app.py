from flask import Flask, render_template, request, jsonify, send_file
from werkzeug.utils import secure_filename
import openai
import os
import tempfile
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = tempfile.mkdtemp()
openai.api_key = "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"  # Thay bằng API key thực

# Xử lý phát hiện ngôn ngữ
def detect_language(text):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "Bạn là một hệ thống phát hiện ngôn ngữ chuyên nghiệp. Hãy phát hiện ngôn ngữ của văn bản sau."},
            {"role": "user", "content": text}
        ]
    )
    return response.choices[0].message['content']

# Xử lý dịch văn bản
def translate_text(text, target_lang):
    source_lang = detect_language(text)
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": f"Bạn là một hệ thống dịch thuật chuyên nghiệp. Hãy dịch chính xác từ {source_lang} sang {target_lang}"},
            {"role": "user", "content": text}
        ]
    )
    return response.choices[0].message['content']

# Routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/translate', methods=['POST'])
def translate():
    data = request.json
    translated = translate_text(data['text'], data['target_lang'])
    return jsonify({'translation': translated})

@app.route('/translate-file', methods=['POST'])
def translate_file():
    file = request.files['file']
    target_lang = request.form['target_lang']
    
    # Xử lý file
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    # Trích xuất văn bản
    text = ""
    if filename.endswith('.pdf'):
        reader = PdfReader(filepath)
        for page in reader.pages:
            text += page.extract_text()
    elif filename.endswith('.docx'):
        doc = Document(filepath)
        text = '\n'.join([para.text for para in doc.paragraphs])
    elif filename.endswith('.pptx'):
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
    
    # Dịch và trả về
    translated = translate_text(text, target_lang)
    return jsonify({'translation': translated})

if __name__ == '__main__':
    app.run(debug=True)